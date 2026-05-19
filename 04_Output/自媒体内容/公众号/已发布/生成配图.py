#!/usr/bin/env python3
"""
公众号文章配图 - 赛博霓虹风
为《为什么你的AI总是装样子》生成配图
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
import os

# 配色方案
NEON_PURPLE = RGBColor(157, 78, 221)
NEON_PINK = RGBColor(255, 20, 147)
NEON_CYAN = RGBColor(0, 255, 255)
NEON_GREEN = RGBColor(57, 255, 20)
NEON_ORANGE = RGBColor(255, 165, 0)
DARK_BG = RGBColor(13, 27, 42)

def set_bg(slide):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = DARK_BG

def add_title(text, slide, size=44, color=NEON_CYAN):
    title = slide.shapes.title
    title.text = text
    title.text_frame.paragraphs[0].font.color.rgb = color
    title.text_frame.paragraphs[0].font.bold = True
    title.text_frame.paragraphs[0].font.size = Pt(size)

def create_cover():
    """头图（封面）"""
    prs = Presentation()
    prs.slide_width = Inches(6)
    prs.slide_height = Inches(3.33)  # 900x500 = 1.8:1
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    set_bg(slide)

    # 装饰线
    for i in range(5):
        shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0.8 + i*0.5), Inches(6), Inches(0.03))
        shape.fill.solid()
        shape.fill.fore_color.rgb = [NEON_PINK, NEON_CYAN, NEON_PURPLE, NEON_GREEN, NEON_ORANGE][i]
        shape.line.fill.background()

    # 主标题
    title = slide.shapes.title
    title.text = "为什么你的 AI 总是「装样子」?"
    title.text_frame.paragraphs[0].font.color.rgb = NEON_CYAN
    title.text_frame.paragraphs[0].font.bold = True
    title.text_frame.paragraphs[0].font.size = Pt(40)

    # 副标题
    sub = slide.placeholders[1]
    sub.text = "因为你缺了这层底座"
    sub.text_frame.paragraphs[0].font.color.rgb = NEON_PINK
    sub.text_frame.paragraphs[0].font.size = Pt(24)

    return prs

def create_pain_point():
    """开头插图 - 痛点"""
    prs = Presentation()
    prs.slide_width = Inches(6)
    prs.slide_height = Inches(3.33)
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    set_bg(slide)

    add_title("你的AI也这样？", slide, 36, NEON_ORANGE)

    pains = [
        "❌ 回答很对，但跟你没关系",
        "❌ 好像专业，但不是你的话",
        "❌ 看起来周到，但用不上"
    ]

    for i, p in enumerate(pains):
        box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(1.5 + i*0.6), Inches(5), Inches(0.5))
        box.fill.solid()
        box.fill.fore_color.rgb = NEON_ORANGE
        box.fill.fore_color.alpha = 40
        box.line.color.rgb = NEON_ORANGE
        tf = box.text_frame
        ptf = tf.paragraphs[0]
        ptf.text = p
        ptf.font.color.rgb = RGBColor(255, 255, 255)
        ptf.font.size = Pt(16)

    return prs

def create_architecture():
    """中间插图 - 三层架构"""
    prs = Presentation()
    prs.slide_width = Inches(6)
    prs.slide_height = Inches(3.33)
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    set_bg(slide)

    add_title("LLM Wiki 三层架构", slide, 32, NEON_CYAN)

    layers = [
        ("Schema 层\n规则/契约", NEON_PINK),
        ("Wiki 层\nAI整理的知识", NEON_CYAN),
        ("Raw Sources\n原始资料", NEON_PURPLE)
    ]

    for i, (text, color) in enumerate(layers):
        box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5 + i*1.8), Inches(1.8), Inches(1.6), Inches(1))
        box.fill.solid()
        box.fill.fore_color.rgb = color
        box.fill.fore_color.alpha = 50
        box.line.color.rgb = color
        tf = box.text_frame
        ptf = tf.paragraphs[0]
        ptf.text = text
        ptf.font.color.rgb = RGBColor(255, 255, 255)
        ptf.font.size = Pt(14)
        ptf.alignment = PP_ALIGN.CENTER

    return prs

def create_ending():
    """结尾插图 - 金句"""
    prs = Presentation()
    prs.slide_width = Inches(6)
    prs.slide_height = Inches(3.33)
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    set_bg(slide)

    # 金句
    box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(1), Inches(5), Inches(1.5))
    box.fill.solid()
    box.fill.fore_color.rgb = NEON_PINK
    box.fill.fore_color.alpha = 30
    box.line.color.rgb = NEON_PINK
    tf = box.text_frame
    ptf = tf.paragraphs[0]
    ptf.text = "真正让人记住的，不是你的专业\n而是你作为一个「活人」的温度和立场"
    ptf.font.color.rgb = RGBColor(255, 255, 255)
    ptf.font.size = Pt(18)
    ptf.alignment = PP_ALIGN.CENTER

    # 标签
    tag = slide.shapes.add_textbox(Inches(0.5), Inches(2.8), Inches(5), Inches(0.4))
    tf = tag.text_frame
    ptf = tf.paragraphs[0]
    ptf.text = "— 蓉蓉子的星球Planet"
    ptf.font.color.rgb = NEON_CYAN
    ptf.font.size = Pt(14)
    ptf.alignment = PP_ALIGN.CENTER

    return prs

def save_ppts():
    output_dir = "/Users/rong/Desktop/ob大脑/04_Output/自媒体内容/公众号/已发布"
    os.makedirs(output_dir, exist_ok=True)

    prs = create_cover()
    prs.save(f"{output_dir}/配图1_头图.pptx")
    print("✅ 配图1_头图.pptx")

    prs = create_pain_point()
    prs.save(f"{output_dir}/配图2_痛点.pptx")
    print("✅ 配图2_痛点.pptx")

    prs = create_architecture()
    prs.save(f"{output_dir}/配图3_三层架构.pptx")
    print("✅ 配图3_三层架构.pptx")

    prs = create_ending()
    prs.save(f"{output_dir}/配图4_金句.pptx")
    print("✅ 配图4_金句.pptx")

if __name__ == "__main__":
    save_ppts()
    print("\n🎨 4张配图已生成，可以导入 enhanced-publisher 发布到公众号")