#!/usr/bin/env python3
"""
直播1 PPT - 六间房霓虹风格
基于 python-pptx 创建
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# 霓虹配色方案
NEON_PURPLE = RGBColor(157, 78, 221)
NEON_PINK = RGBColor(255, 20, 147)
NEON_CYAN = RGBColor(0, 255, 255)
NEON_GREEN = RGBColor(57, 255, 20)
NEON_ORANGE = RGBColor(255, 165, 0)
DARK_BG = RGBColor(10, 10, 20)
DARK_BG_LIGHT = RGBColor(20, 20, 40)

def set_dark_background(slide):
    """设置深色背景"""
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = DARK_BG

def add_title_text(slide, text, size=44, color=NEON_CYAN, top=Inches(0.5)):
    """添加标题文本"""
    title = slide.shapes.add_textbox(Inches(0.5), top, Inches(12.333), Inches(1))
    tf = title.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.color.rgb = color
    p.font.bold = True
    p.font.size = Pt(size)
    p.alignment = PP_ALIGN.CENTER

def add_content_text(slide, text, top=Inches(2), left=Inches(0.5), width=Inches(9), height=Inches(5),
                   font_size=18, color=RGBColor(220, 220, 220)):
    """添加正文文本"""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True

    for i, line in enumerate(text.split('\n')):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = line
        p.font.color.rgb = color
        p.font.size = Pt(font_size)
        p.space_after = Pt(10)

def add_neon_box(slide, text, x, y, w, h, color=NEON_PURPLE):
    """添加霓虹色矩形框"""
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.fill.fore_color.alpha = 50
    shape.line.color.rgb = color
    shape.line.width = Pt(2)

    tf = shape.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.color.rgb = RGBColor(255, 255, 255)
    p.font.size = Pt(14)
    p.alignment = PP_ALIGN.CENTER

def create_ppt():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # Slide 1: 封面
    slide = prs.slides.add_slide(prs.slide_layouts[0])  # Title Slide
    set_dark_background(slide)

    # 装饰线条
    for i in range(6):
        shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(1.5 + i*0.8), Inches(13.333), Inches(0.05))
        shape.fill.solid()
        shape.fill.fore_color.rgb = [NEON_PURPLE, NEON_PINK, NEON_CYAN, NEON_GREEN, NEON_ORANGE, NEON_CYAN][i]
        shape.line.fill.background()

    # 标题
    title = slide.shapes.title
    title.text = "直播 1｜认知升级 × 三层架构 × 跑通 MVP"
    title.text_frame.paragraphs[0].font.color.rgb = NEON_CYAN
    title.text_frame.paragraphs[0].font.bold = True
    title.text_frame.paragraphs[0].font.size = Pt(48)

    # 副标题
    sub = slide.placeholders[1]
    sub.text = "21天 AI 实操营 · 六间房 LLM Wiki"
    sub.text_frame.paragraphs[0].font.color.rgb = NEON_PINK
    sub.text_frame.paragraphs[0].font.size = Pt(24)

    # Slide 2: 目录
    slide = prs.slides.add_slide(prs.slide_layouts[5])  # Blank
    set_dark_background(slide)
    add_title_text(slide, "📋 直播 1 课程结构")

    items = [
        "Part 0 · 开场暖场（5分钟）",
        "Part 1 · 认知升级：为什么装资料 ≠ 知识库（25分钟）",
        "Part 2 · 现场搭建：三层目录 + CLAUDE.md（35分钟）",
        "Part 3 · 跑通第一个 MVP（25分钟）",
        "Part 4 · Q&A + 作业 + 搭子（10分钟）"
    ]

    for i, item in enumerate(items):
        add_neon_box(slide, item, Inches(1), Inches(1.8 + i*0.9), Inches(11.333), Inches(0.7),
                    [NEON_CYAN, NEON_PURPLE, NEON_PINK, NEON_GREEN, NEON_ORANGE][i])

    # Slide 3: 痛点共鸣
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    set_dark_background(slide)
    add_title_text(slide, "🤔 你有没有同感？", size=44, color=NEON_PINK)

    pain_points = [
        "❌ Notion / 飞书 / Obsidian 都试过，但没一个用得起来",
        "❌ 收藏了几百篇好文章，但写东西时还是从零开始",
        "❌ 用 ChatGPT/Claude 时，它给你的回答总像在敷衍"
    ]

    for i, p in enumerate(pain_points):
        add_neon_box(slide, p, Inches(1), Inches(2 + i*1.2), Inches(11.333), Inches(0.8), NEON_ORANGE)

    add_content_text(slide, "💡 你不是不努力，你是**用错了范式**。\n你做的是『资料堆积』，但这个时代真正有用的是**让 AI 读懂你的大脑**。",
                 top=Inches(5.5), font_size=20, color=NEON_CYAN)

    # Slide 4: Karpathy LLM Wiki 三层架构
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    set_dark_background(slide)
    add_title_text(slide, "🧠 Karpathy LLM Wiki 三层架构")

    # 三个层次矩形
    layers = [
        ("Schema 层（规则/契约）\n你维护", NEON_PINK),
        ("Wiki 层（已编译的知识）\nAI 写，你读", NEON_CYAN),
        ("Raw Sources 层（原始资料）\n你写，AI 只读", NEON_PURPLE)
    ]

    for i, (text, color) in enumerate(layers):
        add_neon_box(slide, text, Inches(1 + i*4), Inches(2.5), Inches(3.5), Inches(2), color)

    # 箭头
    for i in range(2):
        arrow = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, Inches(4.3 + i*4), Inches(3.3), Inches(1), Inches(0.5))
        arrow.fill.solid()
        arrow.fill.fore_color.rgb = NEON_GREEN
        arrow.line.fill.background()

    add_content_text(slide, "📌 传统 RAG = 每次去图书馆翻书\nLLM Wiki = 你雇了一个私人图书管理员",
                 top=Inches(5.2), font_size=16, color=RGBColor(180, 180, 180))

    # Slide 5: RAG vs LLM Wiki 对比
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    set_dark_background(slide)
    add_title_text(slide, "🆚 RAG vs LLM Wiki")

    # 表格对比
    table = slide.shapes.add_table(6, 3, Inches(0.5), Inches(1.8), Inches(12.333), Inches(4)).table

    headers = ["维度", "传统 RAG", "LLM Wiki"]
    data = [
        ["知识处理时机", "查询时处理（每次重复）", "摄入时处理（一次性编译）"],
        ["交叉引用", "临时发现", "预先构建并持续维护"],
        ["矛盾检测", "容易被忽略", "摄入时主动标记"],
        ["知识积累", "无—每次从头开始", "复利增长"],
        ["类比", "每次去图书馆翻书", "你有一个私人图书管理员"]
    ]

    # 设置表头
    for i, header in enumerate(headers):
        cell = table.cell(0, i)
        cell.text = header
        cell.fill.solid()
        cell.fill.fore_color.rgb = NEON_PURPLE
        cell.text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
        cell.text_frame.paragraphs[0].font.bold = True

    # 设置数据
    for row_idx, row_data in enumerate(data):
        for col_idx, cell_data in enumerate(row_data):
            cell = table.cell(row_idx + 1, col_idx)
            cell.text = cell_data
            cell.fill.solid()
            if row_idx % 2 == 0:
                cell.fill.fore_color.rgb = DARK_BG_LIGHT
            else:
                cell.fill.fore_color.rgb = DARK_BG
            cell.text_frame.paragraphs[0].font.color.rgb = RGBColor(220, 220, 220)

    # Slide 6: 六间房 × LLM Wiki 联动
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    set_dark_background(slide)
    add_title_text(slide, "🏠 六间房 × LLM Wiki")

    rooms = [
        ("01_我是谁", NEON_CYAN),
        ("02_我会什么", NEON_PURPLE),
        ("03_我说过什么", NEON_PINK),
        ("04_我卖什么", NEON_ORANGE),
        ("05_我面对谁", NEON_GREEN),
        ("06_我怎么想", NEON_PINK)
    ]

    for i, (room, color) in enumerate(rooms):
        add_neon_box(slide, room, Inches(0.5 + (i%3)*4.2), Inches(2 + (i//3)*1.2), Inches(3.8), Inches(0.9), color)

    add_content_text(slide, "💡 六间房解决「装什么」，LLM Wiki 解决「怎么活」\nCLAUDE.md 是你和 AI 之间的契约",
                     top=Inches(5), font_size=20, color=NEON_CYAN)

    # Slide 7: 三个误区
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    set_dark_background(slide)
    add_title_text(slide, "⚠️ 三个误区（避坑指南）")

    misconceptions = [
        ("误区1", "「先把所有资料整理好，再用 AI」", "❌ 资料永远整理不完。\n正确：搭好框架 → 边丢边整理"),
        ("误区2", "「我得先学完 Obsidian 所有功能」", "❌ Obsidian 只是文件夹查看器。\n正确：会用文件夹和 Markdown 就行"),
        ("误区3", "「AI 输出不像我，是提示词没写好」", "❌ 根因是缺少 Schema + Wiki。\n提示词只能解决 5% 的问题")
    ]

    for i, (title, subtitle, content) in enumerate(misconceptions):
        add_neon_box(slide, f"{title}\n{subtitle}", Inches(0.5 + i*4.2), Inches(2), Inches(3.8), Inches(2.2), NEON_ORANGE if i == 0 else NEON_PINK if i == 1 else NEON_CYAN)
        content_box = slide.shapes.add_textbox(Inches(0.5 + i*4.2), Inches(4.4), Inches(3.8), Inches(1.5))
        tf = content_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = content
        p.font.color.rgb = RGBColor(200, 200, 200)
        p.font.size = Pt(12)

    # Slide 8: 目录结构
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    set_dark_background(slide)
    add_title_text(slide, "🛠️ 目录结构搭建", size=44, color=NEON_GREEN)

    code_text = '''ob-大脑/
├── raw/                    ← 你的素材，AI只读
│   ├── chats/              ← 微信对话/语音转文字
│   ├── articles/           ← 收藏的好文章
│   ├── notes/              ← 灵感碎片
│   └── assets/             ← 图片/截图
├── wiki/                   ← AI帮你写
│   ├── 01-identity-os/     ← 我是谁
│   ├── 02-expertise/       ← 我会什么
│   ├── 03-content-lib/    ← 我说过什么
│   ├── 04-products/       ← 我卖什么
│   ├── 05-audience/       ← 我面对谁
│   ├── 06-decision-engine/ ← 我怎么想
│   ├── index.md            ← 全局索引
│   └── log.md              ← 操作日志
└── CLAUDE.md               ← Schema 契约'''

    code_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.8), Inches(12.333), Inches(5))
    tf = code_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = code_text
    p.font.color.rgb = NEON_GREEN
    p.font.size = Pt(11)
    p.font.name = "Consolas"

    # Slide 9: CLAUDE.md 模板
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    set_dark_background(slide)
    add_title_text(slide, "📝 CLAUDE.md 核心��块")

    sections = [
        ("一", "我是谁（一句话定位）"),
        ("二", "六间房含义（AI 用这个分类素材）"),
        ("三", "写作风格（AI 模仿这个语气）"),
        ("四", "母题清单（AI 写内容时按这个分）"),
        ("五", "Ingest 规则（AI 整理素材的方法）"),
        ("六", "Query 规则（AI 回答问题的方法）"),
        ("七", "边界规则（AI 必须守的红线）")
    ]

    for i, (num, text) in enumerate(sections):
        add_neon_box(slide, f"第{num}节：{text}", Inches(0.5 + (i%2)*6.5), Inches(1.8 + (i//2)*0.85), Inches(6), Inches(0.7),
                    [NEON_CYAN, NEON_PINK, NEON_PURPLE, NEON_GREEN, NEON_ORANGE, NEON_PINK, NEON_CYAN][i])

    add_content_text(slide, "💡 DAY 1 只需填前 4 节，后面的先复制模板",
                   top=Inches(5.5), font_size=20, color=NEON_ORANGE)

    # Slide 10: 谁写谁读边界
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    set_dark_background(slide)
    add_title_text(slide, "🚧 谁写谁读边界", size=44, color=NEON_ORANGE)

    boundary_table = slide.shapes.add_table(5, 3, Inches(1), Inches(2), Inches(11.333), Inches(3.5)).table

    headers = ["文件��", "你写？", "AI 写？"]
    data = [
        ["raw/", "✅ 只能在这里加东西", "❌ 不许动"],
        ["wiki/", "❌ 不要直接改", "✅ AI 维护，你只读"],
        ["CLAUDE.md", "✅ 你和 AI 一起改", "✅"],
        ["wiki/log.md", "❌", "✅ AI 追加操作日志"]
    ]

    # 表头
    for i, header in enumerate(headers):
        cell = boundary_table.cell(0, i)
        cell.text = header
        cell.fill.solid()
        cell.fill.fore_color.rgb = NEON_PINK
        cell.text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
        cell.text_frame.paragraphs[0].font.bold = True

    # 数据
    for row_idx, row_data in enumerate(data):
        for col_idx, cell_data in enumerate(row_data):
            cell = boundary_table.cell(row_idx + 1, col_idx)
            cell.text = cell_data
            cell.fill.solid()
            cell.fill.fore_color.rgb = DARK_BG_LIGHT
            cell.text_frame.paragraphs[0].font.color.rgb = RGBColor(220, 220, 220)

    add_content_text(slide, "💡 raw 是档案室，wiki 是图书馆，CLAUDE.md 是契约",
                   top=Inches(6), font_size=18, color=NEON_CYAN)

    # Slide 11: index.md/log.md
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    set_dark_background(slide)
    add_title_text(slide, "🗂️ Wiki 的脊梁：index.md + log.md")

    files = [
        ("index.md", "门牌号系统", "AI 每次 Query 必读，从这里定位「该去哪个房间」"),
        ("log.md", "操作账本", "记录每次 Ingest/Query/Lint 操作，append-only，永不修改")
    ]

    for i, (name, role, desc) in enumerate(files):
        add_neon_box(slide, f"📄 {name}\n{role}", Inches(1 + i*5), Inches(2), Inches(4.5), Inches(1.8),
                    NEON_CYAN if i == 0 else NEON_PINK)

        desc_box = slide.shapes.add_textbox(Inches(1 + i*5), Inches(4), Inches(4.5), Inches(1.5))
        tf = desc_box.text_frame
        p = tf.paragraphs[0]
        p.text = desc
        p.font.color.rgb = RGBColor(200, 200, 200)
        p.font.size = Pt(14)

    # Slide 12: 手写 vs AI 提炼
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    set_dark_background(slide)
    add_title_text(slide, "✋ 手写 vs 🤖 AI 提炼", size=44, color=NEON_PINK)

    # 手动部分
    manual_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(2), Inches(6), Inches(4))
    manual_box.fill.solid()
    manual_box.fill.fore_color.rgb = NEON_PINK
    manual_box.fill.fore_color.alpha = 30
    manual_box.line.color.rgb = NEON_PINK
    manual_box.line.width = Pt(2)

    tf = manual_box.text_frame
    p = tf.paragraphs[0]
    p.text = "✋ 必须手写（3 张地基卡）\n━━━━━━━━━━━━━━━━━━━\n①01_我是谁/一句话定位\n②06_我怎么想/核心价值观\n③CLAUDE.md 写作风格"
    p.font.color.rgb = RGBColor(255, 255, 255)
    p.font.size = Pt(14)
    p.alignment = PP_ALIGN.CENTER

    # AI 部分
    ai_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6.833), Inches(2), Inches(6), Inches(4))
    ai_box.fill.solid()
    ai_box.fill.fore_color.rgb = NEON_CYAN
    ai_box.fill.fore_color.alpha = 30
    ai_box.line.color.rgb = NEON_CYAN
    ai_box.line.width = Pt(2)

    tf = ai_box.text_frame
    p = tf.paragraphs[0]
    p.text = "🤖 默认 AI 提炼（80%）\n━━━━━━━━━━━━━━━━━━━\n你：把素材丢进 raw/\nAI：提取要点 → 写信息卡 → 更新 index\n\n典型：02/03/04/05/案例库"
    p.font.color.rgb = RGBColor(255, 255, 255)
    p.font.size = Pt(14)
    p.alignment = PP_ALIGN.CENTER

    add_content_text(slide, "💡 你做「被调用率最高」的���，AI 做重复性的事",
                     top=Inches(6.3), font_size=18, color=NEON_GREEN)

    # Slide 13: Part 3 跑通 MVP
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    set_dark_background(slide)
    add_title_text(slide, "🚀 Part 3：跑通第一个 MVP", size=44, color=NEON_GREEN)

    steps = [
        ("3.1 ✋", "填 01_我是谁 第一张信息卡", "≈8分钟"),
        ("3.2 🤖", "第一次 Ingest：丢素材让 AI 整理", "≈10分钟"),
        ("3.3 🔍", "第一次 Query：向你的 Wiki 提问", "≈7分钟")
    ]

    for i, (step, title, time) in enumerate(steps):
        add_neon_box(slide, f"{step}\n{title}\n⏱️ {time}", Inches(0.5 + i*4.2), Inches(2), Inches(3.8), Inches(2.2),
                    [NEON_CYAN, NEON_PINK, NEON_PURPLE][i])

    add_content_text(slide, "💡 这就是 LLM Wiki 的复利—你今天填的卡，AI 明天能用、一年后还能用",
                   top=Inches(5), font_size=18, color=NEON_GREEN)

    # Slide 14: 作业清单
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    set_dark_background(slide)
    add_title_text(slide, "📋 Day 1-7 作业清单")

    tasks = [
        "✅ Day 1：建三层目录 + CLAUDE.md 初版 + 一句话定位 + 跑通 Ingest/Query",
        "✅ Day 2-3：丰富 02（把方法论文章丢给 AI 提炼）",
        "✅ Day 4-5：丰富 03 + 05（公众号文章 + 咨询记录）",
        "✅ Day 6：手写 06 核心价值观（地基卡）",
        "✅ Day 7：里程碑验收（5 间房各 ≥1 张卡）"
    ]

    for i, task in enumerate(tasks):
        add_neon_box(slide, task, Inches(0.5), Inches(1.8 + i*1), Inches(12.333), Inches(0.85),
                    [NEON_CYAN, NEON_PINK, NEON_PURPLE, NEON_GREEN, NEON_ORANGE][i])

    # Slide 15: Q&A
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    set_dark_background(slide)
    add_title_text(slide, "❓ Q&A")

    qas = [
        ("Q1", "Obsidian 打不开？", "换网页版"),
        ("Q2", "CLAUDE.md 写不出？", "先用模板，跑两天再调"),
        ("Q3", "04_我卖什么要填？", "不必，先空着"),
        ("Q4", "Wiki 写错了？", "让 AI 改，别直接改"),
        ("Q5", "raw 能放 PDF？", "可以，建议先转 MD")
    ]

    for i, (q, question, answer) in enumerate(qas):
        add_neon_box(slide, f"{q}: {question}\n→ {answer}", Inches(0.5 + (i%2)*6.5), Inches(2 + (i//2)*1.3), Inches(6), Inches(1),
                    NEON_PINK if i % 2 == 0 else NEON_CYAN)

    # Slide 16: 结束页
    slide = prs.slides.add_slide(prs.slide_layouts[0])  # Title Slide
    set_dark_background(slide)

    # 装饰线
    for i in range(6):
        shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(1.5 + i*0.8), Inches(13.333), Inches(0.05))
        shape.fill.solid()
        shape.fill.fore_color.rgb = [NEON_PINK, NEON_CYAN, NEON_PURPLE, NEON_GREEN, NEON_ORANGE, NEON_CYAN][i]
        shape.line.fill.background()

    title = slide.shapes.title
    title.text = "21 天后，你有一只会自己长大的 AI 助手"
    title.text_frame.paragraphs[0].font.color.rgb = NEON_CYAN
    title.text_frame.paragraphs[0].font.bold = True
    title.text_frame.paragraphs[0].font.size = Pt(42)

    sub = slide.placeholders[1]
    sub.text = "六间房 LLM Wiki · 实操营"
    sub.text_frame.paragraphs[0].font.color.rgb = NEON_PINK
    sub.text_frame.paragraphs[0].font.size = Pt(24)

    # 保存文件
    output_path = "/Users/rong/Desktop/ob大脑/04_Output/课件资料/2026-05_实操营升级版/直播1_PPT_霓虹风格.pptx"
    prs.save(output_path)
    print(f"✅ PPT 已保存到: {output_path}")

if __name__ == "__main__":
    create_ppt()